import io, uuid, math
from typing import Dict, Any, List, Optional
from datetime import datetime

from docx import Document
from docx.shared import Inches, Pt, Cm, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT
from docx.oxml.ns import qn

from reportlab.lib.pagesizes import A4
from reportlab.lib import colors
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak, Image as RLImage

from pptx import Presentation
from pptx.util import Inches as PptInches, Pt as PptPt
from pptx.dml.color import RGBColor as PptRGB
from pptx.enum.text import PP_ALIGN

from app.services.analytics.analysis_service import AnalysisService
from app.services.analytics.student_tracking import StudentTrackingService, ClassAnalysisService
from app.services.report.chart_generator import draw_bar_chart, draw_line_chart, draw_radar_chart, draw_table_image

# Register Chinese font for PDF
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
import os as _os
_CJK_REGISTERED = False
_CJK_FONT_CANDIDATES = [
    ('C:/Windows/Fonts/simsun.ttc', True),
    ('C:/Windows/Fonts/msyh.ttc', True),
    ('C:/Windows/Fonts/simhei.ttf', True),
    ('/System/Library/Fonts/PingFang.ttc', False),
]
for _path, _is_win in _CJK_FONT_CANDIDATES:
    if _os.path.exists(_path):
        try:
            pdfmetrics.registerFont(TTFont("CJK-Font", _path))
            _CJK_REGISTERED = True
            break
        except:
            continue

class ReportService:
    def _get_cjk_style(self, name, parent, **kw):
        """Create a paragraph style with CJK font support"""
        from reportlab.lib.styles import ParagraphStyle
        try:
            style = ParagraphStyle(name, parent=parent, **kw)
            style.fontName = "CJK-Font"
        except:
            style = ParagraphStyle(name, parent=parent, **kw)
        return style

    def __init__(self, db):
        self.db = db
        self.analysis_service = AnalysisService(db)

    def get_analysis_data(self, exam_id: str) -> Dict[str, Any]:
        return self.analysis_service.get_exam_analysis(exam_id)


    # ==================== CLASS ANALYSIS ====================

    def _add_chart_to_docx(self, doc, chart_buf, width_inches=5.5):
        """Add a chart image to a Word document"""
        if chart_buf and chart_buf.getbuffer().nbytes > 1000:
            doc.add_picture(chart_buf, width=Inches(width_inches))
            doc.add_paragraph("")

    def _add_chart_to_pdf(self, elements, chart_buf, width=460, height=260):
        """Add a chart image to PDF elements list"""
        import sys as _sys
        if chart_buf:
            buf_data = chart_buf.getvalue()
            print(f"DEBUG: chart buf size: {len(buf_data)}", file=_sys.stderr)
            if len(buf_data) > 500:
                chart_buf.seek(0)
                try:
                    elements.append(RLImage(chart_buf, width=width, height=height))
                    elements.append(Spacer(1, 8))
                except Exception as e:
                    print(f"DEBUG: chart add error: {e}", file=_sys.stderr)

    def generate_class_word(self, class_id: str) -> io.BytesIO:
        from app.services.analytics.student_tracking import ClassAnalysisService
        service = ClassAnalysisService(self.db)
        data = service.get_class_overview(class_id)
        if not data:
            return None

        doc = Document()
        style = doc.styles['Normal']
        font = style.font
        font.name = 'SimSun'
        font.size = Pt(11)
        style.element.rPr.rFonts.set(qn('w:eastAsia'), 'SimSun')

        # Title
        title = doc.add_heading(f"{data.get('class_name', '')} 班级分析报告", level=0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        doc.add_paragraph(f"年级: {data.get('grade_name', '-')}    学生人数: {data.get('student_count', 0)}    考试次数: {data.get('exam_count', 0)}")
        doc.add_paragraph("")

        # Overview table
        doc.add_heading("各科统计", level=1)
        subject_stats = data.get("subject_stats", [])
        if subject_stats:
            table = doc.add_table(rows=1, cols=5)
            table.style = 'Light Grid Accent 1'
            table.alignment = WD_TABLE_ALIGNMENT.CENTER
            hdr = table.rows[0].cells
            for i, h in enumerate(["科目", "平均分", "最高分", "最低分", "样本数"]):
                hdr[i].text = h
            for ss in subject_stats:
                row = table.add_row().cells
                row[0].text = ss.get("subject_name", "")
                row[1].text = str(ss.get("avg_score", ""))
                row[2].text = str(ss.get("max_score", ""))
                row[3].text = str(ss.get("min_score", ""))
                row[4].text = str(ss.get("count", ""))

        # Bar chart - subject averages
        if subject_stats:
            labels = [s.get("subject_name", "") for s in subject_stats]
            scores = [s.get("avg_score", 0) for s in subject_stats]
            chart = draw_bar_chart(scores, labels, title="各科平均分", ylabel="平均分")
            self._add_chart_to_docx(doc, chart)

        # Exam trend
        doc.add_heading("考试成绩趋势", level=1)
        exam_summary = data.get("exam_summary", [])
        if exam_summary:
            table = doc.add_table(rows=1, cols=3)
            table.style = 'Light Grid Accent 1'
            for i, h in enumerate(["考试名称", "考试日期", "平均得分率(%)"]):
                table.rows[0].cells[i].text = h
            for es in exam_summary:
                row = table.add_row().cells
                row[0].text = es.get("exam_name", "")
                row[1].text = es.get("exam_date", "")
                row[2].text = f"{es.get('avg_rate', '')}%"

            # Line chart
            exam_names = [e.get("exam_name", "")[:8] for e in exam_summary]
            rates = [e.get("avg_rate", 0) or 0 for e in exam_summary]
            if rates:
                chart = draw_line_chart([{"name": "得分率", "data": rates}], exam_names, title="考试成绩趋势", ylabel="得分率(%)")
                self._add_chart_to_docx(doc, chart)

        # Student list
        doc.add_heading("学生成绩一览", level=1)
        student_list = data.get("student_list", [])
        if student_list:
            table = doc.add_table(rows=1, cols=3)
            table.style = 'Light Grid Accent 1'
            for i, h in enumerate(["学号", "姓名", "综合得分率(%)"]):
                table.rows[0].cells[i].text = h
            for sl in student_list:
                row = table.add_row().cells
                row[0].text = sl.get("student_no", "")
                row[1].text = sl.get("student_name", "")
                row[2].text = f"{sl.get('avg_rate', '')}%" if sl.get('avg_rate') is not None else "-"

        # Risk students
        risk = data.get("risk_students", [])
        if risk:
            doc.add_heading("需关注学生", level=1)
            table = doc.add_table(rows=1, cols=3)
            table.style = 'Light Grid Accent 1'
            for i, h in enumerate(["学号", "姓名", "综合得分率(%)"]):
                table.rows[0].cells[i].text = h
            for r in risk:
                row = table.add_row().cells
                row[0].text = r.get("student_no", "")
                row[1].text = r.get("student_name", "")
                row[2].text = f"{r.get('avg_rate', '')}%"

        buf = io.BytesIO()
        doc.save(buf)
        buf.seek(0)
        return buf

    def generate_class_pdf(self, class_id: str) -> io.BytesIO:
        from app.services.analytics.student_tracking import ClassAnalysisService
        service = ClassAnalysisService(self.db)
        data = service.get_class_overview(class_id)
        if not data:
            return None

        buf = io.BytesIO()
        doc = SimpleDocTemplate(buf, pagesize=A4)
        styles = getSampleStyleSheet()
        title_style = self._get_cjk_style('CustomTitle', styles['Title'], fontSize=18, spaceAfter=12)
        heading_style = self._get_cjk_style('CustomHeading', styles['Heading2'], fontSize=14, spaceBefore=10, spaceAfter=6)
        normal = self._get_cjk_style('CJKNormal', styles['Normal'])

        elements = []
        elements.append(Paragraph(f"{data.get('class_name', '')} 班级分析报告", title_style))
        elements.append(Paragraph(f"年级: {data.get('grade_name', '-')}    学生人数: {data.get('student_count', 0)}    考试次数: {data.get('exam_count', 0)}", normal))
        elements.append(Spacer(1, 12))

        # Subject stats table
        elements.append(Paragraph("各科统计", heading_style))
        subject_stats = data.get("subject_stats", [])
        if subject_stats:
            t_data = [[s.get("subject_name", ""), str(s.get("avg_score", "")), str(s.get("max_score", "")), str(s.get("min_score", "")), str(s.get("count", ""))] for s in subject_stats]
            t = Table([["科目", "平均分", "最高分", "最低分", "样本数"]] + t_data, colWidths=[80, 70, 70, 70, 60])
            t.setStyle(TableStyle([
                ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#409EFF')),
                ('TEXTCOLOR', (0, 0), (-1, 0), colors.white),
                ('FONTNAME', (0, 0), (-1, -1), 'CJK-Font'),
                ('FONTSIZE', (0, 0), (-1, -1), 9),
                ('GRID', (0, 0), (-1, -1), 0.5, colors.grey),
                ('ALIGN', (1, 1), (-1, -1), 'CENTER'),
            ]))
            elements.append(t)
            elements.append(Spacer(1, 8))

        # Subject bar chart
        if subject_stats:
            try:
                labels = [s.get("subject_name", "") for s in subject_stats]
                scores = [s.get("avg_score", 0) for s in subject_stats]
                chart_buf = draw_bar_chart(scores, labels, title="各科平均分", ylabel="平均分", width=520, height=320)
                if chart_buf and len(chart_buf.getvalue()) > 500:
                    from reportlab.platypus import Image as RLImage
                    chart_buf.seek(0)
                    elements.append(PageBreak())
                    elements.append(Paragraph("各科平均分柱状图", heading_style))
                    elements.append(RLImage(chart_buf, width=460, height=280))
                    elements.append(Spacer(1, 12))
            except:
                pass

        elements.append(Paragraph("考试成绩趋势", heading_style))
        exam_summary = data.get("exam_summary", [])
        if exam_summary:
            t_data = [[e.get("exam_name", ""), e.get("exam_date", ""), f"{e.get('avg_rate', '')}%"] for e in exam_summary]
            t = Table([["考试名称", "考试日期", "平均得分率(%)"]] + t_data, colWidths=[160, 80, 100])
            t.setStyle(TableStyle([
                ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#409EFF')),
                ('TEXTCOLOR', (0, 0), (-1, 0), colors.white),
                ('FONTNAME', (0, 0), (-1, -1), 'CJK-Font'),
                ('FONTSIZE', (0, 0), (-1, -1), 9),
                ('GRID', (0, 0), (-1, -1), 0.5, colors.grey),
                ('ALIGN', (1, 1), (-1, -1), 'CENTER'),
            ]))
            elements.append(t)
            elements.append(Spacer(1, 8))

        # Exam trend line chart
        if exam_summary:
            try:
                rates = [e.get("avg_rate", 0) or 0 for e in exam_summary]
                exam_names = [e.get("exam_name", "")[:8] for e in exam_summary]
                if rates:
                    chart_buf = draw_line_chart([{"name": "得分率", "data": rates}], exam_names, title="考试成绩趋势", width=520, height=300)
                    if chart_buf and len(chart_buf.getvalue()) > 500:
                        from reportlab.platypus import Image as RLImage
                        chart_buf.seek(0)
                        elements.append(PageBreak())
                        elements.append(Paragraph("考试成绩趋势图", heading_style))
                        elements.append(RLImage(chart_buf, width=460, height=260))
                        elements.append(Spacer(1, 12))
            except:
                pass

        elements.append(Paragraph("学生成绩一览", heading_style))
        student_list = data.get("student_list", [])
        if student_list:
            t_data = [[s.get("student_no", ""), s.get("student_name", ""), f"{s.get('avg_rate', '')}%" if s.get('avg_rate') is not None else "-"] for s in student_list[:30]]
            t = Table([["学号", "姓名", "得分率(%)"]] + t_data, colWidths=[110, 90, 80])
            t.setStyle(TableStyle([
                ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#409EFF')),
                ('TEXTCOLOR', (0, 0), (-1, 0), colors.white),
                ('FONTNAME', (0, 0), (-1, -1), 'CJK-Font'),
                ('FONTSIZE', (0, 0), (-1, -1), 8),
                ('GRID', (0, 0), (-1, -1), 0.3, colors.grey),
                ('ALIGN', (2, 1), (-1, -1), 'CENTER'),
            ]))
            elements.append(t)
            if len(student_list) > 30:
                elements.append(Paragraph(f"...共 {len(student_list)} 名学生，仅展示前30名", normal))

        doc.build(elements)
        buf.seek(0)
        return buf

    # ==================== STUDENT ANALYSIS ====================

    def generate_student_word(self, student_id: str) -> io.BytesIO:
        service = StudentTrackingService(self.db)
        data = service.get_student_scores(student_id)
        if not data:
            return None

        doc = Document()
        style = doc.styles['Normal']
        font = style.font
        font.name = 'SimSun'
        font.size = Pt(11)
        style.element.rPr.rFonts.set(qn('w:eastAsia'), 'SimSun')

        # Title
        title = doc.add_heading(f"{data.get('student_name', '')} 学情跟踪报告", level=0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        doc.add_paragraph(f"学号: {data.get('student_no', '-')}    班级: {data.get('class_name', '-')}    年级: {data.get('grade_name', '-')}")
        doc.add_paragraph(f"考试次数: {data.get('exam_count', 0)}    总体趋势: {data.get('overall_trend', '暂无数据')}")
        doc.add_paragraph("")

        # Strengths & Weaknesses
        doc.add_heading("优势与薄弱科目", level=1)
        strengths = data.get("strengths", [])
        weaknesses = data.get("weaknesses", [])
        if strengths or weaknesses:
            doc.add_paragraph(f"优势科目 ({len(strengths)}): " + ", ".join([f"{s.get('subject_name', '')}({s.get('avg_rate', '')}%)" for s in strengths]))
            doc.add_paragraph(f"薄弱科目 ({len(weaknesses)}): " + ", ".join([f"{s.get('subject_name', '')}({s.get('avg_rate', '')}%)" for s in weaknesses]))
            doc.add_paragraph("")

        # Trend chart
        doc.add_heading("各科成绩趋势", level=1)
        exams = data.get("exams", [])
        trends = data.get("trends", [])
        if exams and trends:
            exam_names = [e.get("exam_name", "")[:6] for e in exams]
            series_list = []
            for t in trends:
                scores = [(s.get("rate", 0) or 0) for s in (t.get("scores") or [])]
                if scores:
                    series_list.append({"name": t.get("subject_name", ""), "data": scores})
            if series_list and exam_names:
                chart = draw_line_chart(series_list, exam_names, title="各科成绩趋势", ylabel="得分率(%)", width=640, height=360)
                self._add_chart_to_docx(doc, chart)

        # Exam history table
        doc.add_heading("历次成绩", level=1)
        if exams:
            table = doc.add_table(rows=1, cols=3 + max((len(e.get("subjects", [])) for e in exams), default=0))
            table.style = 'Light Grid Accent 1'
            hdr = table.rows[0].cells
            hdr[0].text = "考试"
            hdr[1].text = "日期"
            hdr[2].text = "综合得分率"
            col = 3
            for e in exams:
                for s in (e.get("subjects") or []):
                    if col < len(hdr):
                        hdr[col].text = s.get("subject_name", "")
                    col += 1
                break  # Only need first exam's subjects for headers
            for e in exams:
                row = table.add_row().cells
                row[0].text = e.get("exam_name", "")
                row[1].text = e.get("exam_date", "")
                row[2].text = f"{e.get('avg_rate', '')}%"
                ci = 3
                for s in (e.get("subjects") or []):
                    if ci < len(row):
                        row[ci].text = str(s.get("score", "-"))
                        ci += 1

        buf = io.BytesIO()
        doc.save(buf)
        buf.seek(0)
        return buf

    def generate_student_pdf(self, student_id: str) -> io.BytesIO:
        service = StudentTrackingService(self.db)
        data = service.get_student_scores(student_id)
        if not data:
            return None

        buf = io.BytesIO()
        doc = SimpleDocTemplate(buf, pagesize=A4)
        styles = getSampleStyleSheet()
        title_style = self._get_cjk_style('CustomTitle', styles['Title'], fontSize=18, spaceAfter=12)
        heading_style = self._get_cjk_style('CustomHeading', styles['Heading2'], fontSize=14, spaceBefore=10, spaceAfter=6)
        normal = self._get_cjk_style('CJKNormal', styles['Normal'])
        bold = self._get_cjk_style('BoldNormal', styles['Normal'], fontSize=10, spaceAfter=4)

        elements = []
        elements.append(Paragraph(f"{data.get('student_name', '')} 学情跟踪报告", title_style))
        elements.append(Paragraph(f"学号: {data.get('student_no', '-')}    班级: {data.get('class_name', '-')}    年级: {data.get('grade_name', '-')}", normal))
        elements.append(Paragraph(f"考试次数: {data.get('exam_count', 0)}    总体趋势: {data.get('overall_trend', '暂无数据')}", normal))
        elements.append(Spacer(1, 12))

        # Strengths & Weaknesses
        elements.append(Paragraph("优势与薄弱科目", heading_style))
        strengths = data.get("strengths", [])
        weaknesses = data.get("weaknesses", [])
        if strengths:
            s_text = "; ".join([f"{s.get('subject_name', '')}(平均{s.get('avg_rate', '')}%/最新{s.get('latest_rate', '')}%)" for s in strengths])
            elements.append(Paragraph(f"<b>优势科目:</b> {s_text}", bold))
        if weaknesses:
            w_text = "; ".join([f"{s.get('subject_name', '')}(平均{s.get('avg_rate', '')}%/最新{s.get('latest_rate', '')}%)" for s in weaknesses])
            elements.append(Paragraph(f"<b>薄弱科目:</b> {w_text}", bold))

        # Trend chart
        elements.append(Paragraph("各科成绩趋势", heading_style))
        exams = data.get("exams", [])
        trends = data.get("trends", [])
        if exams and trends:
            exam_names = [e.get("exam_name", "")[:6] for e in exams]
            series_list = []
            for t in trends:
                scores = [(s.get("rate", 0) or 0) for s in (t.get("scores") or [])]
                if scores:
                    series_list.append({"name": t.get("subject_name", ""), "data": scores})
            # Trend chart
        elements.append(Paragraph("各科成绩趋势", heading_style))
        exams = data.get("exams", [])
        trends = data.get("trends", [])
        if exams and trends:
            try:
                exam_names = [e.get("exam_name", "")[:6] for e in exams]
                series_list = []
                for t in trends:
                    scores = [(s.get("rate", 0) or 0) for s in (t.get("scores") or [])]
                    if scores:
                        series_list.append({"name": t.get("subject_name", ""), "data": scores})
                if series_list:
                    chart = draw_line_chart(series_list, exam_names, title="各科成绩趋势", ylabel="得分率(%)", width=520, height=320)
                    if chart and len(chart.getvalue()) > 500:
                        from reportlab.platypus import Image as RLImage
                        chart.seek(0)
                        elements.append(PageBreak())
                        elements.append(Paragraph("各科成绩趋势图", heading_style))
                        elements.append(RLImage(chart, width=460, height=280))
                        elements.append(Spacer(1, 12))
            except:
                pass

        # Exam history
        elements.append(Paragraph("历次成绩", heading_style))
        if exams:
            t_data = []
            for e in exams:
                row = [e.get("exam_name", ""), e.get("exam_date", ""), f"{e.get('avg_rate', '')}%" if e.get('avg_rate') else "-"]
                for s in (e.get("subjects") or []):
                    row.append(f"{s.get('subject_name', '')}:{s.get('score', '-')}")
                t_data.append(row)
            max_cols = max(len(r) for r in t_data) if t_data else 3
            headers = ["考试", "日期", "综合得分率"]
            if t_data:
                extra = max_cols - 3
                for i in range(extra):
                    headers.append("科目" + str(i+1))
            t = Table([headers] + t_data, colWidths=[100, 60, 70] + [50] * (max_cols - 3))
            t.setStyle(TableStyle([
                ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#409EFF')),
                ('TEXTCOLOR', (0, 0), (-1, 0), colors.white),
                ('FONTNAME', (0, 0), (-1, -1), 'CJK-Font'),
                ('FONTSIZE', (0, 0), (-1, -1), 8),
                ('GRID', (0, 0), (-1, -1), 0.3, colors.grey),
                ('ALIGN', (1, 1), (-1, -1), 'CENTER'),
            ]))
            elements.append(t)

        doc.build(elements)
        buf.seek(0)
        return buf

    # ==================== WORD ====================

    def generate_word(self, exam_id: str) -> io.BytesIO:
        from app.services.report.word_generator import generate_exam_word
        return generate_exam_word(exam_id, self.db, self.analysis_service)

    def generate_pdf(self, exam_id: str) -> io.BytesIO:
        data = self.get_analysis_data(exam_id)
        buf = io.BytesIO()
        doc = SimpleDocTemplate(buf, pagesize=A4)

        styles = getSampleStyleSheet()
        title_style = self._get_cjk_style('CustomTitle', styles['Title'], fontSize=18, spaceAfter=20)
        heading_style = self._get_cjk_style('CustomHeading', styles['Heading2'], fontSize=14, spaceBefore=12, spaceAfter=6)
        normal = self._get_cjk_style('CJKNormal', styles['Normal'])

        elements = []
        elements.append(Paragraph(data.get("exam_name", "考试分析报告"), title_style))
        elements.append(Paragraph(f"考试日期: {data.get('exam_date', '-')}", normal))
        elements.append(Paragraph(f"参考人数: {data.get('total_students', 0)}", normal))
        elements.append(Spacer(1, 12))

        # Grade stats
        elements.append(Paragraph("年级总体统计", heading_style))
        grade_stats = data.get("grade_stats", [])
        if grade_stats:
            headers = ["科目", "平均分", "得分率%", "最高分", "最低分", "及格率%", "优秀率%", "标准差"]
            rows = [[
                gs.get("subject_name", ""),
                str(gs.get("avg_score", "")),
                str(gs.get("avg_score_rate", "")),
                str(gs.get("max_score", "")),
                str(gs.get("min_score", "")),
                str(gs.get("pass_rate", "")),
                str(gs.get("excellent_rate", "")),
                str(gs.get("std_dev", ""))
            ] for gs in grade_stats]
            t = Table([headers] + rows, colWidths=[46, 38, 38, 34, 34, 38, 38, 34])
            t.setStyle(TableStyle([
                ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#409EFF')),
                ('TEXTCOLOR', (0, 0), (-1, 0), colors.white),
                ('FONTNAME', (0, 0), (-1, -1), 'CJK-Font'),
                ('FONTSIZE', (0, 0), (-1, -1), 7),
                ('LEADING', (0, 0), (-1, -1), 10),
                ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                ('GRID', (0, 0), (-1, -1), 0.5, colors.grey),
                ('ROWBACKGROUNDS', (0, 1), (-1, -1), [colors.white, colors.HexColor('#f5f7fa')]),
                ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
                ('TOPPADDING', (0, 0), (-1, -1), 3),
                ('BOTTOMPADDING', (0, 0), (-1, -1), 3),
                ('LEFTPADDING', (0, 0), (-1, -1), 2),
                ('RIGHTPADDING', (0, 0), (-1, -1), 2),
            ]))
            elements.append(t)
            elements.append(PageBreak())

        # Class stats
        elements.append(Paragraph("各班统计", heading_style))
        for cs in data.get("class_stats", []):
            elements.append(Paragraph(f"{cs.get('class_name', '')} (人数: {cs.get('student_count', 0)})", heading_style))
            stats_list = cs.get("stats", [])
            if stats_list:
                h = ["科目", "平均分", "得分率%", "及格率%"]
                r = [[s.get("subject_name", ""), str(s.get("avg_score", "")), str(s.get("avg_score_rate", "")), str(s.get("pass_rate", ""))] for s in stats_list]
                t = Table([h] + r, colWidths=[80, 55, 55, 55])
                t.setStyle(TableStyle([
                    ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#409EFF')),
                    ('TEXTCOLOR', (0, 0), (-1, 0), colors.white),
                    ('FONTNAME', (0, 0), (-1, -1), 'CJK-Font'),
                    ('FONTSIZE', (0, 0), (-1, -1), 8),
                    ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                    ('GRID', (0, 0), (-1, -1), 0.5, colors.grey),
                ]))
                elements.append(t)
                elements.append(Spacer(1, 8))

        doc.build(elements)
        buf.seek(0)
        return buf

    # ==================== PPT ====================

    def generate_ppt(self, exam_id: str) -> io.BytesIO:
        data = self.get_analysis_data(exam_id)
        prs = Presentation()
        prs.slide_width = PptInches(13.333)
        prs.slide_height = PptInches(7.5)

        # Color scheme
        BLUE = PptRGB(64, 158, 255)
        DARK = PptRGB(48, 49, 51)
        GRAY = PptRGB(144, 147, 153)
        WHITE = PptRGB(255, 255, 255)

        def add_slide(title_text, subtitle_text=""):
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
            # Title bg
            from pptx.util import Emu
            shape = slide.shapes.add_shape(1, PptInches(0), PptInches(0), prs.slide_width, PptInches(1.2))
            shape.fill.solid()
            shape.fill.fore_color.rgb = BLUE
            shape.line.fill.background()
            tf = shape.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = title_text
            p.font.size = PptPt(28)
            p.font.color.rgb = WHITE
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER

            if subtitle_text:
                p2 = tf.add_paragraph()
                p2.text = subtitle_text
                p2.font.size = PptPt(14)
                p2.font.color.rgb = WHITE
                p2.alignment = PP_ALIGN.CENTER
            return slide

        def add_textbox(slide, left, top, width, height, text, size=14, bold=False, color=DARK, align=PP_ALIGN.LEFT):
            txBox = slide.shapes.add_textbox(PptInches(left), PptInches(top), PptInches(width), PptInches(height))
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = text
            p.font.size = PptPt(size)
            p.font.color.rgb = color
            p.font.bold = bold
            p.alignment = align
            return txBox

        def add_table(slide, headers, rows, left, top, width, height):
            tbl = slide.shapes.add_table(len(rows)+1, len(headers), PptInches(left), PptInches(top), PptInches(width), PptInches(height)).table
            for i, h in enumerate(headers):
                cell = tbl.cell(0, i)
                cell.text = h
                for p in cell.text_frame.paragraphs:
                    p.font.size = PptPt(10)
                    p.font.bold = True
                    p.font.color.rgb = WHITE
                cell.fill.solid()
                cell.fill.fore_color.rgb = BLUE
            for ri, row in enumerate(rows):
                for ci, val in enumerate(row):
                    cell = tbl.cell(ri+1, ci)
                    cell.text = str(val)
                    for p in cell.text_frame.paragraphs:
                        p.font.size = PptPt(10)
            return tbl

        # Slide 1: Cover
        slide1 = add_slide(data.get("exam_name", "考试质量分析报告"), f"考试日期: {data.get('exam_date', '-')}  |  参考人数: {data.get('total_students', 0)}")
        add_textbox(slide1, 1, 3, 11, 1, "考试质量分析系统 - 自动生成报告", size=16, color=GRAY, align=PP_ALIGN.CENTER)

        # Slide 2: Grade stats
        slide2 = add_slide("年级总体成绩统计")
        grade_stats = data.get("grade_stats", [])
        if grade_stats:
            headers = ["科目", "平均分", "得分率", "最高分", "最低分", "及格率", "优秀率"]
            rows = [[
                gs.get("subject_name", ""),
                str(gs.get("avg_score", "")),
                f"{gs.get('avg_score_rate', '')}%",
                str(gs.get("max_score", "")),
                str(gs.get("min_score", "")),
                f"{gs.get('pass_rate', '')}%",
                f"{gs.get('excellent_rate', '')}%"
            ] for gs in grade_stats]
            add_table(slide2, headers, rows, 0.5, 1.5, 12, 4)

        # Slide 3-5: Class stats
        for ci, cs in enumerate(data.get("class_stats", [])):
            slide = add_slide(f"班级分析 - {cs.get('class_name', '')}", f"人数: {cs.get('student_count', 0)}")
            stats_list = cs.get("stats", [])
            if stats_list:
                headers = ["科目", "平均分", "得分率", "及格率", "优秀率"]
                rows = [[
                    s.get("subject_name", ""),
                    str(s.get("avg_score", "")),
                    f"{s.get('avg_score_rate', '')}%",
                    f"{s.get('pass_rate', '')}%",
                    f"{s.get('excellent_rate', '')}%"
                ] for s in stats_list]
                add_table(slide, headers, rows, 0.5, 1.5, 12, 4)
            if ci >= 4:
                break

        # Slide 6: Summary
        slide6 = add_slide("总结")
        summary_lines = []
        if grade_stats:
            best = max(grade_stats, key=lambda x: x.get("avg_score_rate", 0))
            worst = min(grade_stats, key=lambda x: x.get("avg_score_rate", 0))
            summary_lines.append(f"得分率最高科目: {best.get('subject_name', '')} ({best.get('avg_score_rate', '')}%)")
            summary_lines.append(f"得分率最低科目: {worst.get('subject_name', '')} ({worst.get('avg_score_rate', '')}%)")
            summary_lines.append("")
            summary_lines.append("建议:")
            summary_lines.append(f"1. 加强 {worst.get('subject_name', '')} 的教学辅导")
            summary_lines.append("2. 关注低分段学生的学习情况")
            summary_lines.append("3. 保持优势科目的教学力度")
        y = 2
        for line in summary_lines:
            add_textbox(slide6, 1, y, 11, 0.5, line, size=14 if not line.startswith("建议") else 16, bold=line.startswith("建议"))
            y += 0.5

        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        return buf


