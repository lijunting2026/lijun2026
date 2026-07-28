import io
from typing import Dict, Any
from pptx.util import Inches as PptInches, Pt as PptPt
from pptx.dml.color import RGBColor as PptRGB
from pptx.enum.text import PP_ALIGN

from app.services.report.chart_generator import draw_bar_chart, draw_line_chart, draw_radar_chart, draw_table_image

class PptGenerator:
    def __init__(self, db):
        self.db = db

    def generate_ppt(self, exam_id: str) -> io.BytesIO:
        data = self.get_analysis_data(exam_id)
        prs = Presentation()
        prs.slide_width = PptInches(13.333)
        prs.slide_height = PptInches(7.5)

        # Color scheme
        BLUE = PptRGB(64, 158, 255)
        DARK = PptRGB(48, 49, 51)
        GRAY = PptRGB(144, 147, 153)
        WHITE = PptRGB(255, 255, 255)

        def add_slide(title_text, subtitle_text=""):
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
            # Title bg
            from pptx.util import Emu
            shape = slide.shapes.add_shape(1, PptInches(0), PptInches(0), prs.slide_width, PptInches(1.2))
            shape.fill.solid()
            shape.fill.fore_color.rgb = BLUE
            shape.line.fill.background()
            tf = shape.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = title_text
            p.font.size = PptPt(28)
            p.font.color.rgb = WHITE
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER

            if subtitle_text:
                p2 = tf.add_paragraph()
                p2.text = subtitle_text
                p2.font.size = PptPt(14)
                p2.font.color.rgb = WHITE
                p2.alignment = PP_ALIGN.CENTER
            return slide

        def add_textbox(slide, left, top, width, height, text, size=14, bold=False, color=DARK, align=PP_ALIGN.LEFT):
            txBox = slide.shapes.add_textbox(PptInches(left), PptInches(top), PptInches(width), PptInches(height))
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = text
            p.font.size = PptPt(size)
            p.font.color.rgb = color
            p.font.bold = bold
            p.alignment = align
            return txBox

        def add_table(slide, headers, rows, left, top, width, height):
            tbl = slide.shapes.add_table(len(rows)+1, len(headers), PptInches(left), PptInches(top), PptInches(width), PptInches(height)).table
            for i, h in enumerate(headers):
                cell = tbl.cell(0, i)
                cell.text = h
                for p in cell.text_frame.paragraphs:
                    p.font.size = PptPt(10)
                    p.font.bold = True
                    p.font.color.rgb = WHITE
                cell.fill.solid()
                cell.fill.fore_color.rgb = BLUE
            for ri, row in enumerate(rows):
                for ci, val in enumerate(row):
                    cell = tbl.cell(ri+1, ci)
                    cell.text = str(val)
                    for p in cell.text_frame.paragraphs:
                        p.font.size = PptPt(10)
            return tbl

        # Slide 1: Cover
        slide1 = add_slide(data.get("exam_name", "考试质量分析报告"), f"考试日期: {data.get('exam_date', '-')}  |  参考人数: {data.get('total_students', 0)}")
        add_textbox(slide1, 1, 3, 11, 1, "考试质量分析系统 - 自动生成报告", size=16, color=GRAY, align=PP_ALIGN.CENTER)

        # Slide 2: Grade stats
        slide2 = add_slide("年级总体成绩统计")
        grade_stats = data.get("grade_stats", [])
        if grade_stats:
            headers = ["科目", "平均分", "得分率", "最高分", "最低分", "及格率", "优秀率"]
            rows = [[
                gs.get("subject_name", ""),
                str(gs.get("avg_score", "")),
                f"{gs.get('avg_score_rate', '')}%",
                str(gs.get("max_score", "")),
                str(gs.get("min_score", "")),
                f"{gs.get('pass_rate', '')}%",
                f"{gs.get('excellent_rate', '')}%"
            ] for gs in grade_stats]
            add_table(slide2, headers, rows, 0.5, 1.5, 12, 4)

        # Slide 3-5: Class stats
        for ci, cs in enumerate(data.get("class_stats", [])):
            slide = add_slide(f"班级分析 - {cs.get('class_name', '')}", f"人数: {cs.get('student_count', 0)}")
            stats_list = cs.get("stats", [])
            if stats_list:
                headers = ["科目", "平均分", "得分率", "及格率", "优秀率"]
                rows = [[
                    s.get("subject_name", ""),
                    str(s.get("avg_score", "")),
                    f"{s.get('avg_score_rate', '')}%",
                    f"{s.get('pass_rate', '')}%",
                    f"{s.get('excellent_rate', '')}%"
                ] for s in stats_list]
                add_table(slide, headers, rows, 0.5, 1.5, 12, 4)
            if ci >= 4:
                break

        # Slide 6: Summary
        slide6 = add_slide("总结")
        summary_lines = []
        if grade_stats:
            best = max(grade_stats, key=lambda x: x.get("avg_score_rate", 0))
            worst = min(grade_stats, key=lambda x: x.get("avg_score_rate", 0))
            summary_lines.append(f"得分率最高科目: {best.get('subject_name', '')} ({best.get('avg_score_rate', '')}%)")
            summary_lines.append(f"得分率最低科目: {worst.get('subject_name', '')} ({worst.get('avg_score_rate', '')}%)")
            summary_lines.append("")
            summary_lines.append("建议:")
            summary_lines.append(f"1. 加强 {worst.get('subject_name', '')} 的教学辅导")
            summary_lines.append("2. 关注低分段学生的学习情况")
            summary_lines.append("3. 保持优势科目的教学力度")
        y = 2
        for line in summary_lines:
            add_textbox(slide6, 1, y, 11, 0.5, line, size=14 if not line.startswith("建议") else 16, bold=line.startswith("建议"))
            y += 0.5

        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        return buf



